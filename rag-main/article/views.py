from django.shortcuts import render

# 导入数据模型ArticlePost
from .models import ArticlePost
# 导入 HttpResponse 模块
from django.http import HttpResponse

# 视图函数
def article_home(request):
    # 取出所有博客文章
    articles = ArticlePost.objects.all()
    # 需要传递给模板（templates）的对象
    context = {'articles': articles}
    # render函数：载入模板，并返回context对象
    return render(request, 'article/home.html', context)




#chat_with_gpt视图函数，用于与gpt沟通，augment_prompt函数是在接入数据库之后，
#把问题和数据库中前k个相似度最高的文档组合起来。注意其中的API 密钥和 API 基础地址，
#我是使用的第三方api，所以需要改基础地址。
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import json
from langchain.schema import SystemMessage, HumanMessage
from langchain.chat_models import ChatOpenAI

def augment_prompt(query: str, k: int):#augment_prompt函数是在接入数据库之后，把问题和数据库中前k个相似度最高的文档组合起来
    results = vectorstore.similarity_search(query, k=k)
    source_knowledge = "\n".join([x.page_content for x in results])
    augmented_prompt = f"""基于以下内容并且使用中文回答问题：
    
内容:
{source_knowledge}

query: {query}"""
    return augmented_prompt

@csrf_exempt
def chat_with_gpt(request):#视图函数
    if request.method == 'POST':
        try:
            # 打印请求体调试信息
            print("Request body:", request.body)

            # 解析 JSON 数据
            data = json.loads(request.body)
            message = data.get('message')

            # 读取 context.txt 文件内容
            if os.path.exists(CONTEXT_FILE_PATH):
                with open(CONTEXT_FILE_PATH, 'r', encoding='utf-8') as f:
                    saved_context = f.read()
            else:
                saved_context = ""

            # 合并用户输入和上下文
            message = f"{saved_context}\n用户问题: {message}"

            # 你的 API 密钥和 API 基础地址
            openai_api_key = "sk-qXWtQCRRjUJG637NAf1e347317794d84Ad264c9f10B3779a"
            openai_api_base = "https://www.jcapikey.com/v1"

            chat = ChatOpenAI(
                openai_api_key=openai_api_key,
                openai_api_base=openai_api_base,
                model='gpt-3.5-turbo'
            )
            #组合问题和前k个最相似文本
            augmented_prompt = augment_prompt(message, 10)
            messages = [
                SystemMessage(content="你是一个专业的知识助手。"),
                HumanMessage(content=augmented_prompt)
            ]
            #与gpt对话
            response = chat(messages)
            gpt_response = response.content.strip()  # 修正响应内容
            return JsonResponse({'response': gpt_response})

        except json.JSONDecodeError:
            return JsonResponse({'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            # 打印异常信息
            print("Exception:", str(e))
            return JsonResponse({'error': str(e)}, status=400)

    return JsonResponse({'error': 'Invalid request method'}, status=405)



#upload_file视图函数，用于处理上传的文件，并且把处理完的数据整合到数据库中
import os
from django.conf import settings
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.core.files.storage import default_storage
from langchain.document_loaders import PyPDFLoader, CSVLoader,UnstructuredHTMLLoader,Docx2txtLoader,UnstructuredFileLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import Chroma
from langchain.schema import SystemMessage, HumanMessage
from langchain.chat_models import ChatOpenAI
from langchain.schema import Document
from utils.xlsx import process_excel_files
from pptx import Presentation

# 设置文件上传目录和向量库存储目录，上传的文件保存在/media/uploads下，
#保存的向量库存储在/media/vectorstore下。
UPLOAD_FOLDER = os.path.join(settings.MEDIA_ROOT, 'uploads')
VECTORSTORE_PATH = os.path.join(settings.MEDIA_ROOT, 'vectorstore')
VECTORSTORE_METADATA_PATH = os.path.join(VECTORSTORE_PATH, "metadata.txt")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(VECTORSTORE_PATH, exist_ok=True)

# 替换为你的 API 密钥和 API 基础地址
openai_api_key = "sk-qXWtQCRRjUJG637NAf1e347317794d84Ad264c9f10B3779a"
openai_api_base = "https://www.jcapikey.com/v1"

#关于ppt和pptx的处理方法
def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text_runs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return '\n'.join(text_runs)

def process_file(file_path, file_type):#检查文件类型，对于不同文件的处理方法
    #读取文件
    if file_type == 'pdf':
        loader = PyPDFLoader(file_path)
        pages = loader.load_and_split()
    elif file_type == 'csv':
        loader = CSVLoader(file_path=file_path)
        pages = loader.load()
    elif file_type == 'html':
        loader = UnstructuredHTMLLoader(file_path)
        pages = loader.load_and_split()
    elif file_type == 'docx':
        loader = Docx2txtLoader(file_path=file_path)
        pages = loader.load()
    elif file_type == 'doc':
        loader = Docx2txtLoader(file_path=file_path)
        pages = loader.load()
    elif file_type == 'xlsx':
        # 调用新的Excel处理函数
        excel_data = process_excel_files(file_path)
        pages = [Document(page_content="\n".join([" ".join(map(str, row)) for row in sheet['content']])) for sheet in excel_data]
    elif file_type == 'txt':
        loader = UnstructuredFileLoader(file_path)
        text=loader.load()
        doc = Document(page_content=text, metadata={"source": file_path})
        pages = [doc]
    elif file_type == 'pptx':
        text = extract_text_from_pptx(file_path)
        doc = Document(page_content=text, metadata={"source": file_path})
        pages = [doc]
    #文件切割
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50,
    )
    docs = text_splitter.split_documents(pages)
    #使用embedding预训练模型，处理数据，最后返回向量库
    embed_model = OpenAIEmbeddings(openai_api_key=openai_api_key, openai_api_base=openai_api_base)
    vectorstore = Chroma.from_documents(documents=docs, embedding=embed_model, persist_directory=VECTORSTORE_PATH)
    return vectorstore

# def save_vectorstore(vectorstore):#保存向量库
#     vectorstore.persist()

def save_vectorstore(vectorstore, file_names):#保存向量库和已储存文件
    vectorstore.persist()
    with open(VECTORSTORE_METADATA_PATH, 'a') as f:
        for name in file_names:
            f.write(name + '\n')

def load_vectorstore_metadata():#读取向量库存储文件名
    if os.path.exists(VECTORSTORE_METADATA_PATH):
        with open(VECTORSTORE_METADATA_PATH, 'r') as f:
            return f.read().splitlines()
    return []

def load_vectorstore():#读取向量库
    if os.path.exists(VECTORSTORE_PATH):
        return Chroma(persist_directory=VECTORSTORE_PATH, embedding_function=OpenAIEmbeddings(openai_api_key=openai_api_key, openai_api_base=openai_api_base))
    else:
        return None

vectorstore = load_vectorstore()

def allowed_file(filename):#检测上传文档中的目标类型文档
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'pdf', 'csv', 'html','doc','docx','xlsx','txt','pptx'}

@csrf_exempt
def upload_file(request):
    name_list = []
    text = '成功上传文件：'
    global vectorstore

    if request.method == 'POST':
        files = request.FILES.getlist('files[]')  # 使用正确的键来获取文件列表
        if not files:
            return JsonResponse({'error': 'No file part'})

        for file in files:
            if file.name == '':
                return JsonResponse({'error': 'No selected file'})

            if file and allowed_file(file.name):
                filename = os.path.basename(file.name)  # 使用 secure_filename 来确保文件名安全
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                with default_storage.open(file_path, 'wb+') as destination:
                    for chunk in file.chunks():
                        destination.write(chunk)

                file_type = filename.rsplit('.', 1)[1].lower()
                vectorstore = process_file(file_path, file_type)
                name_list.append(file.name)
            else:
                return JsonResponse({'error': 'Invalid file type'})

        save_vectorstore(vectorstore, name_list)

        for i in name_list:
            text += str(i) + " "

        return JsonResponse({'message': text})
    else:
        return JsonResponse({'error': 'Invalid request method'}, status=405)


@csrf_exempt #视图函数，获取向量库文件名
def get_vectorstore_metadata(request):
    if request.method == 'GET':
        files = load_vectorstore_metadata()
        return JsonResponse({'files': files})
    else:
        return JsonResponse({'error': 'Invalid request method'}, status=405)


# 定义保存上下文的路径
CONTEXT_FILE_PATH = os.path.join(settings.MEDIA_ROOT, 'context.txt')

@csrf_exempt
def save_context(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            context = data.get('context')

            if not context:
                return JsonResponse({'error': 'Context cannot be empty'}, status=400)

            # 将新内容写入到 context.txt 文件中
            with open(CONTEXT_FILE_PATH, 'w', encoding='utf-8') as f:
                f.write(context)

            return JsonResponse({'message': 'Context successfully updated'})

        except json.JSONDecodeError:
            return JsonResponse({'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=400)

    return JsonResponse({'error': 'Invalid request method'}, status=405)

